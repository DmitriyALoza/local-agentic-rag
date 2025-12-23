"""
PowerPoint (PPTX) Document Parser
Extracts text from slides, titles, notes, and shapes
"""

from pathlib import Path
from typing import List, Dict
from pptx import Presentation


class PowerPointParser:
    """Parser for PowerPoint (.pptx) files"""

    def parse(self, file_path: str) -> List[Dict]:
        """
        Parse a PowerPoint file and extract structured content

        Args:
            file_path: Path to the .pptx file

        Returns:
            List of dictionaries, each containing:
            - text: The extracted text content
            - section: Section identifier (e.g., "Slide 1", "Slide 2 Notes")
            - metadata: Additional metadata about the content
        """
        prs = Presentation(file_path)
        chunks = []
        filename = Path(file_path).name

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Extract slide title
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text

            # Extract text from all shapes
            slide_text_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)

            # Combine all text from the slide
            slide_text = "\n".join(slide_text_parts)

            # Create chunk for slide content
            if slide_text.strip():
                chunks.append({
                    "text": slide_text,
                    "section": f"Slide {slide_num}" + (f": {title_text}" if title_text else ""),
                    "metadata": {
                        "filename": filename,
                        "slide_number": slide_num,
                        "has_title": bool(title_text),
                        "title": title_text
                    }
                })

            # Extract notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    chunks.append({
                        "text": notes_text,
                        "section": f"Slide {slide_num} Notes",
                        "metadata": {
                            "filename": filename,
                            "slide_number": slide_num,
                            "content_type": "notes",
                            "related_title": title_text
                        }
                    })

        return chunks

    def get_metadata(self, file_path: str) -> Dict:
        """
        Get metadata about the PowerPoint file

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary containing file metadata
        """
        prs = Presentation(file_path)
        path = Path(file_path)

        return {
            "filename": path.name,
            "file_type": "pptx",
            "total_slides": len(prs.slides),
            "file_size_bytes": path.stat().st_size
        }
